from pptx import Presentation
import pptx

ppt1 = Presentation("original.pptx")
ppt2 = Presentation("new.pptx")

def ppt_to_text(ppt):

    initial_string = ""
    presentation_array = []
    to_write = ""
    
    # Redunant export to text files, used for testing only #
    # ppt_data = open(f"text_ppt_{id(ppt)}.txt", "a")

    for slide in ppt.slides:

        slide_number = ppt.slides.index(slide)
        initial_string += f"\nSlide #{str(slide_number)}"
        # ppt_data.write(f"\nSlide #{str(slide_number)}")

        
        for shape in slide.shapes:

            # If images are found, they are converted to binary strings for comparison #
            if isinstance(shape, pptx.shapes.picture.Picture):  
                blob = shape.image.blob
                # ppt_data.write(f"<----- Image binarystring -----> \n{str(blob)} \n\n") 
                initial_string += f"<----- Image binarystring -----> \n{str(blob)} \n\n"

            if not shape.has_text_frame:
                continue

            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    # ppt_data.write(run.text)
                    to_write += run.text

                initial_string += f" {to_write} "
        
        presentation_array.append(initial_string)
        initial_string = ""
    
    return presentation_array

original_presentation = ppt_to_text(ppt1)
new_presentation = ppt_to_text(ppt2)

def compare_presentations(original, new):


    discrepancies = 0
    f = open("results.txt", "a")

    if len(original) != len(new):
        f.write()
        return

    for i in range (0, len(new)):
        
        if original[i] == new[i]:
            pass
        else:
            discrepancies += 1
            f.write(f"Discrepancies have been found at Slide #{i + 1}\n")
    
    if discrepancies == 0:
        f.write(f"Discrepancies have not been found! Congratulations.")
    f.close()

compare_presentations(original_presentation, new_presentation)